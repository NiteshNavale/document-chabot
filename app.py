import streamlit as st
from dotenv import load_dotenv
import os
import pandas as pd
from io import StringIO
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser

# --- PAGE CONFIGURATION ---
st.set_page_config(
    page_title="NitBot",
    page_icon="🤖",
    layout="wide"
)

# --- ENVIRONMENT AND API KEY LOADING ---
load_dotenv()
groq_api_key = os.getenv("GROQ_API_KEY")

# --- UI STYLING (CSS) ---
st.markdown("""
<style>
    /* --- Main App Styling --- */
    .stApp {
        background-color: #e5ddd5; /* WhatsApp-like fallback color */
        background-image: radial-gradient(circle, rgba(0,0,0,0.05) 1px, transparent 1px);
        background-size: 15px 15px; /* Adjust size to make the pattern more/less dense */
        font-family: 'Helvetica Neue', sans-serif;
    }

    /* --- Sidebar Styling --- */
    .st-emotion-cache-16txtl3 {
        background-color: #F8F9FA;
        border-right: 1px solid #D1D7DB;
    }

    /* --- Create the Chat Window "Card" --- */
    .st-emotion-cache-1y4p8pa {
        padding: 2rem 3rem;
        background-color: #F0F2F6;
        border-radius: 1rem;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
        margin: 1rem;
    }

    /* --- Chat Header Styling --- */
    .st-emotion-cache-1yyk08v {
        background-color: #FFFFFF;
        padding: 1rem;
        border-radius: 10px 10px 0 0;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        border-bottom: 1px solid #E0E0E0;
        text-align: center;
        font-weight: 600;
        color: #1E1E1E;
    }

    /* --- WhatsApp-like Chat Bubbles --- */
    [data-testid="stChatMessages"] {
        display: flex;
        flex-direction: column;
        padding: 1rem 0;
    }

    /* General chat message container styling */
    [data-testid="stChatMessage"] {
        border-radius: 12px;
        padding: 10px 15px;
        margin-bottom: 0.75rem;
        box-shadow: 0 1px 2px 0 rgba(0,0,0,0.15);
        border: none;
    }

    /* User message styling */
    [data-testid="stChatMessage"]:has([data-testid="stAvatarIcon-user"]) {
        background-color: #dcf8c6;
        align-self: flex-end;
        margin-left: auto;
        width: fit-content;
        max-width: 65%;
    }

    /* Bot message styling */
    [data-testid="stChatMessage"]:has([data-testid="stAvatarIcon-assistant"]) {
        background-color: #ffffff;
        align-self: flex-start;
        margin-right: auto;
        width: 100%;
        max-width: 100%;
    }

    /* Remove the avatar icons */
    [data-testid="stAvatarIcon-user"], [data-testid="stAvatarIcon-assistant"] {
        display: none;
    }
    
    /* Table styling */
    .stDataFrame, .stTable {
        border-radius: 8px;
        box-shadow: none;
        border: 1px solid #E0E0E0;
    }

    /* Title styling */
    h1 {
        text-align: center;
        color: #4A4A4A;
        padding-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)

# --- CORE FUNCTIONS ---

# UPDATED: Added logic to read .xlsx files using pandas
def get_docs_text(docs):
    text = ""
    for doc in docs:
        try:
            doc.seek(0)
            file_extension = os.path.splitext(doc.name)[1].lower()
            if file_extension == '.pdf':
                pdf_reader = PdfReader(doc)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text: text += page_text + "\n"
            elif file_extension == '.docx':
                document = docx.Document(doc)
                for para in document.paragraphs: text += para.text + "\n"
            elif file_extension == '.pptx':
                prs = Presentation(doc)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): text += shape.text + "\n"
            elif file_extension == '.txt':
                text += doc.getvalue().decode("utf-8", errors='ignore') + "\n"
            
            # --- NEW: Added logic for Excel files ---
            elif file_extension in ['.xlsx', '.xls']:
                try:
                    # Use pandas ExcelFile to access all sheets
                    xls = pd.ExcelFile(doc, engine='openpyxl')
                    for sheet_name in xls.sheet_names:
                        # Add sheet name as a header for context
                        text += f"\n--- Content of sheet: {sheet_name} ---\n\n"
                        # Read each sheet into a DataFrame
                        df = pd.read_excel(xls, sheet_name=sheet_name)
                        # Convert DataFrame to a string and append
                        text += df.to_string(index=False) + "\n\n"
                except Exception as e:
                    st.warning(f"Could not parse Excel file '{doc.name}'. Error: {e}")
                    # Fallback to reading as raw text if parsing fails
                    doc.seek(0)
                    text += doc.getvalue().decode("utf-8", errors='ignore') + "\n"
            # --- End of new logic ---
            
            elif file_extension == '.csv':
                try:
                    common_kwargs = {'on_bad_lines': 'skip', 'engine': 'python'}
                    doc.seek(0)
                    try:
                        df = pd.read_csv(doc, encoding='utf-8', **common_kwargs)
                    except UnicodeDecodeError:
                        doc.seek(0)
                        df = pd.read_csv(doc, encoding='latin-1', **common_kwargs)
                    text += df.to_string(index=False) + "\n\n"
                except Exception as e:
                    st.warning(f"Could not parse CSV '{doc.name}' as a table. Reading as raw text. Error: {e}")
                    doc.seek(0)
                    text += doc.getvalue().decode("utf-8", errors='ignore') + "\n"
        except Exception as e:
            st.error(f"An unexpected error occurred while processing {doc.name}: {e}")
            continue
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    try:
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
        return True
    except Exception as e:
        st.error(f"Error creating vector store: {e}")
        return False

def get_conversational_chain():
    prompt_template = """
    You are an expert assistant. You will be given a question and a set of context extracted from a document.
    Synthesize this information into a single, coherent final response. Your goal is to be as helpful as possible.
    Base your answer ONLY on the context provided.
    CRITICAL INSTRUCTION: If the user's question asks for a comparison, a list of items, a summary of features,
    or any other structured data that would benefit from a tabular format, YOUR FINAL ANSWER MUST be formatted as a Markdown table.
    Use columns and rows appropriately. Do not just list items; structure them in a table.
    For all other questions, provide a clear, well-formatted text answer.
    CONTEXT:
    {context}
    QUESTION:
    {question}
    Final Answer:
    """
    model = ChatGroq(model_name="llama3-8b-8192", temperature=0.2, api_key=groq_api_key)
    prompt = PromptTemplate.from_template(prompt_template)
    return prompt | model | StrOutputParser()

# --- MAIN APP LAYOUT ---

def main():
    st.title("🤖 NitBot: Your Intelligent Document Assistant")
    
    # --- SESSION STATE INITIALIZATION ---
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "vector_store" not in st.session_state:
        st.session_state.vector_store = None

    # --- SIDEBAR ---
    with st.sidebar:
        st.header("🛠️ Setup Panel")
        if not groq_api_key:
            st.error("GROQ_API_KEY not found. Please set it in your environment.")
            st.stop()
        
        st.subheader("1. Upload Your Documents")
        # UPDATED: Added 'xlsx' to the list of accepted file types
        uploaded_docs = st.file_uploader("Supports PDF, DOCX, PPTX, TXT, CSV, XLSX", accept_multiple_files=True, type=['pdf', 'docx', 'pptx', 'txt', 'csv', 'xlsx', 'xls'])

        st.subheader("2. Process Documents")
        if st.button("Process", use_container_width=True):
            if uploaded_docs:
                with st.spinner("Reading, chunking, and embedding... please wait."):
                    raw_text = get_docs_text(uploaded_docs)
                    if not raw_text.strip():
                        st.warning("No text extracted. Check document content or format.")
                    else:
                        text_chunks = get_text_chunks(raw_text)
                        if get_vector_store(text_chunks):
                            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
                            st.session_state.vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
                            st.session_state.chat_history = [] 
                            st.success("Knowledge base is ready!")
                            st.rerun() 
            else:
                st.warning("Please upload at least one document.")

    # --- MAIN CHAT INTERFACE ---
    st.header("💬 Chat with NitBot")
    if st.session_state.vector_store is None:
        st.info("Please process your documents in the sidebar to begin the chat.")

    # Display chat history from session state
    for message in st.session_state.chat_history:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    # Handle user input and streaming response
    if user_question := st.chat_input("Ask a question about your documents..."):
        if st.session_state.vector_store is None:
            st.warning("Please upload and process documents before asking a question.")
            st.stop()
        
        st.session_state.chat_history.append({"role": "user", "content": user_question})
        with st.chat_message("user"):
            st.markdown(user_question)

        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                try:
                    retriever = st.session_state.vector_store.as_retriever(search_type="similarity", search_kwargs={"k": 4})
                    docs = retriever.invoke(user_question)
                    context = "\n\n---\n\n".join([doc.page_content for doc in docs])
                    chain = get_conversational_chain()
                    full_response = st.write_stream(chain.stream({
                        "context": context,
                        "question": user_question
                    }))
                    st.session_state.chat_history.append({"role": "assistant", "content": full_response})
                except Exception as e:
                    error_message = f"An error occurred: {e}"
                    st.error(error_message)
                    st.session_state.chat_history.append({"role": "assistant", "content": error_message})

if __name__ == "__main__":
    main()
